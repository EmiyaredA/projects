from pptx import Presentation
import os
from pptx.util import Inches
import copy
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO

def set_content(slide, text, type):
    match type:
        case "title":
            title = slide.shapes.title
            title.text = text
        case "content":
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = text
        case _:
            raise ValueError("设置ppt内容阶段，slide中元素的类型传参错误")
def generate_ppt(template_path):
    prs = Presentation(template_path)

    # text_runs 将会被填充为一个字符串列表，
    # 每个字符串代表演示文稿中的一段文字运行
    text_runs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)

    pass


def add_one_slide(save_path):
    SLD_LAYOUT_TITLE_AND_CONTENT = 1

    prs = Presentation()
    slide_layout = prs.slide_layouts[SLD_LAYOUT_TITLE_AND_CONTENT]
    slide = prs.slides.add_slide(slide_layout)
    prs.save(save_path)

def search_placeholder():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    # 对 placeholders 集合的项访问类似于字典而不是列表。虽然上面使用的键是一个整数，但查找的是 idx 值，而不是序列中的位置。
    # 如果提供的值与其中一个占位符的 idx 值不匹配，则 KeyError 将被引发。IDX 值不一定是连续的。
    for shape in slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))


def identify_and_characterize_placeholder():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    for shape in slide.shapes:
        print('%s' % shape.shape_type)

def check_placeholder_format():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    for shape in slide.shapes:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('%d, %s' % (phf.idx, phf.type))

def insert_image_into_placeholder(image_route):
    """PicturePlaceholder.insert_picture()"""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])
    placeholder = slide.placeholders[1]
    print(placeholder.name)
    print(placeholder.placeholder_format.type)
    picture = placeholder.insert_picture(image_route)

def insert_table_into_placeholder(ppt_route):
    """TablePlaceholder.insert_table()"""
    prs = Presentation(ppt_route)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    # example假设一个名为having - table - placeholder.pptx第二个幻灯片版式上具有idx 10的表格占位符
    placeholder = slide.placeholders[10]# idx key, not position
    print(placeholder.name)
    print(placeholder.placeholder_format.type)
    graphic_frame = placeholder.insert_table(rows=2, cols=2)
    table = graphic_frame.table
    print((len(table.rows), len(table.columns)))

def insert_chart_into_placeholder(ppt_route):
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    """ChartPlaceholder.insert_chart()"""
    prs = Presentation(ppt_route)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    placeholder = slide.placeholders[10]  # idx key, not position
    print(placeholder.name)
    print(placeholder.placeholder_format.type)
    chart_data = ChartData()
    chart_data.categories = ['Yes', 'No']
    chart_data.add_series('Series 1', (42, 24))
    # 请注意，insert_chart（） 的返回值是 PlaceholderGraphicFrame 对象，而不是图表本身。
    # PlaceholderGraphicFrame 对象具有 GraphicFrame 对象的所有属性和方法，
    # 以及特定于占位符的属性和方法。插入的图表包含在图形框架中，可以使用其 chart 属性获取。
    graphic_frame = placeholder.insert_chart(XL_CHART_TYPE.PIE, chart_data)
    chart = graphic_frame.chart
    print(chart.chart_type)

def set_slide_title(slide, title):
    """Slide.shapes.title"""
    title_placeholder = slide.shapes.title
    title_placeholder.text = title or 'Air-speed Velocity of Unladen Swallows'
def set_slide_content(slide, content):
    """Slide.shapes.placeholders[1]"""
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = content


def access_text_frame(slide):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text_frame = shape.text_frame
def add_text_paragraphs(shape):
    """文本框架始终至少包含一个段落。这会导致将多个段落变成一个形状的过程比人们想要的要笨拙一些。例如，假设你想要一个包含三个段落的形状："""
    paragraph_strs = [
        'Egg, bacon, sausage and spam.',
        'Spam, bacon, sausage and spam.',
        'Spam, egg, spam, spam, bacon and spam.'
    ]

    text_frame = shape.text_frame
    text_frame.clear()  # remove any existing paragraphs, leaving one empty one

    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]

    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str


def add_text(shape):
    """实际上，只有运行可以包含文本。为 .text 分配字符串 属性是用于放置 文本。以下两个代码段生成 相同的结果："""
    shape.text = 'foobar'

    # is equivalent to ...

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'foobar'

def applying_text_frame_level_formatting(shape):
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
    """
    下面的选项将生成一个形状，其中包含单个段落、底部略宽于上边距
    （默认为 0.05 英寸）、无左边距、文本顶部对齐和关闭自动换行。
    此外，自动调整大小行为设置为调整形状的宽度和高度以适合其文本。
    请注意，在文本框架上设置了垂直对齐方式。在每个段落上设置水平对齐方式：
    :return:
    """
    text_frame = shape.text_frame
    text_frame.text = 'Spam, eggs, and spam' # 单个段落
    text_frame.margin_bottom = Inches(0.08) # 底部略宽于上边距
    text_frame.margin_left = 0 # 无左边距
    text_frame.vertical_anchor = MSO_ANCHOR.TOP # 文本顶部对齐
    text_frame.word_wrap = False # 关闭自动换行
    text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT # 调整形状的宽度和高度以适合其文本

def apply_paragraph_formatting(shape):
    from pptx.enum.text import PP_ALIGN
    """
    下面生成一个包含三个左对齐段落的形状，第二个和第三个段落缩进（类似于子项目符号）在第一个段落下
    :param shape:
    :return:
    """
    paragraph_strs = [
        'Egg, bacon, sausage and spam.',
        'Spam, bacon, sausage and spam.',
        'Spam, egg, spam, spam, bacon and spam.'
    ]

    text_frame = shape.text_frame
    text_frame.clear()

    p = text_frame.paragraphs[0]
    p.text = paragraph_strs[0]
    p.alignment = PP_ALIGN.LEFT # 左对齐

    for para_str in paragraph_strs[1:]:
        p = text_frame.add_paragraph()
        p.text = para_str
        p.alignment = PP_ALIGN.LEFT # 左对齐
        p.level = 1


def apply_character_formatting(shape):
    """
    字符级格式化是通过运行级别应用的，使用 `.font` 属性。
    以下代码将句子格式化为 18pt 的 Calibri 粗体，并应用主题颜色 Accent 1。
    :param shape:
    :return:
    """
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.util import Pt

    text_frame = shape.text_frame
    text_frame.clear()  # not necessary for newly-created shape

    p = text_frame.paragraphs[0]
    # 在 PowerPoint 中，run 是指一段具有相同格式的文本。
    # 在一个段落中，文本可以被分成多个 run，每个 run 可以有不同的格式（如字体、颜色、大小等）。
    # 例如，一个段落中的某些文字可以是加粗的，而其他部分则是正常字体。
    run = p.add_run()
    run.text = 'Spam, eggs, and spam'

    font = run.font
    font.name = 'Calibri'
    font.size = Pt(18)
    font.bold = True
    font.italic = None  # cause value to be inherited from theme
    font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    # 如果您愿意，可以将字体颜色设置为绝对 RGB 值。请注意，当主题更改时，这不会改变颜色：
    font.color.rgb = RGBColor(0xFF, 0x7F, 0x50)
    # 还可以通过提供目标 URL 将运行创建为超链接：
    run.hyperlink.address = 'https://github.com/scanny/python-pptx'


def copy_slide(prs, slide):
    # 创建新的幻灯片
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        if shape.is_placeholder:
            continue

        # 复制文本框
        if shape.shape_type == 1:  # 文本框
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            textbox = new_slide.shapes.add_textbox(left, top, width, height)
            textbox.text = shape.text

            # 复制文本格式
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    new_run = textbox.text_frame.paragraphs[0].add_run()
                    new_run.text = run.text
                    new_run.font.size = run.font.size
                    new_run.font.bold = run.font.bold
                    new_run.font.italic = run.font.italic
                    new_run.font.underline = run.font.underline
                    new_run.font.color.rgb = run.font.color.rgb

        # 复制图片
        elif shape.shape_type == 13:  # 图片
            left = shape.left
            top = shape.top
            image_blob = shape.image.blob  # 获取图片的二进制数据
            image_stream = BytesIO(image_blob)  # 将二进制数据转换为流
            new_slide.shapes.add_picture(image_stream, left, top)

        # 复制表格
        elif shape.shape_type == 19:  # 表格
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            table = new_slide.shapes.add_table(shape.table.rows, shape.table.columns, left, top, width, height).table

            for row_idx, row in enumerate(shape.table.rows):
                for col_idx, cell in enumerate(row.cells):
                    new_cell = table.cell(row_idx, col_idx)
                    new_cell.text = cell.text
                    # 复制单元格格式（字体、对齐等）
                    new_cell.text_frame.paragraphs[0].font.size = cell.text_frame.paragraphs[0].font.size
                    new_cell.text_frame.paragraphs[0].font.bold = cell.text_frame.paragraphs[0].font.bold
                    new_cell.text_frame.paragraphs[0].font.italic = cell.text_frame.paragraphs[0].font.italic

        # 复制形状
        elif shape.shape_type == 6:  # 形状（如矩形）
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, left, top, width, height)
            # 复制形状的填充色、边框等
            new_shape.fill.solid()
            new_shape.fill.fore_color.rgb = shape.fill.fore_color.rgb
            new_shape.line.color.rgb = shape.line.color.rgb

        # 复制图表
        elif shape.shape_type == 4:  # 图表
            # 目前python-pptx不直接支持复制图表，我们可以提取图表数据并重新创建图表
            # 这部分需要根据具体需求来实现
            pass

    return new_slide


def delete_slide(template_path, slide_index, save_path):
    presentation = Presentation(template_path)
    xml_slides = presentation.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[slide_index])
    # 保存
    presentation.save(save_path)

def copy_slide_trial(template_path, save_path):
    prs = Presentation(template_path)
    # 创建新的PPT文件
    new_ppt = Presentation()
    # 复制每一页幻灯片
    for slide in prs.slides:
        copy_slide(new_ppt, slide)
    new_ppt.save(save_path)

#todo 后续还有关于tables和chart的使用样例要构建，等使用到了再说

if __name__ == '__main__':
    # 获取当前文件所在的路径
    dir_path = os.path.dirname(os.path.realpath(__file__))
    template_route = os.path.join(dir_path, "Template", "最近PPT最新进展工作汇报.pptx")
    save_route = os.path.join(dir_path, "GeneratedPresentations", "test.pptx")

    generate_ppt(template_route)
    # add_one_slide(save_route)
    # search_placeholder()
    # identify_and_characterize_placeholder()
    # check_placeholder_format()

    # insert_table_into_placeholder()
    # copy_slide_trial(template_path=template_route, save_path=save_route)
    # delete_slide(template_path=template_route, slide_index=0, save_path=save_route)
