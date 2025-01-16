from pptx import Presentation
import os
import re

def copy_slide(prs, slide):
    slide_layout = slide.slide_layout  # 使用原幻灯片的布局
    new_slide = prs.slides.add_slide(slide_layout)  # 创建新幻灯片

    # 复制原幻灯片上的所有形状
    for shape in slide.shapes:
        if shape.is_placeholder:
            # 对占位符类型的形状进行复制
            new_shape = new_slide.shapes.add_placeholder(shape.placeholder_format.idx, shape.left, shape.top,
                                                         shape.width, shape.height)
            new_shape.text = shape.text  # 复制文本内容（如果有）
        elif shape.shape_type == 13:  # 图片（图像类型为13）
            # 复制图片
            new_slide.shapes.add_picture(shape.image, shape.left, shape.top, shape.width, shape.height)
        else:
            # 其他形状（如矩形、圆形等）
            new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width,
                                                   shape.height)
            new_shape.text = shape.text  # 复制文本内容（如果有）

    return new_slide

def fill_ppt_content(template_route, save_route):

    def match_placeholder(tag, text):
        pattern =  repr(tag) + r"(\d+)"

        # 获取第一个匹配项
        match = re.search(pattern, text.lower())

        if match:
            return int(match.group(1))  # 输出第一个匹配的数字内容
        return None

    # 打开PPT文件
    prs = Presentation(template_route)
    # 遍历所有幻灯片
    xml_slides = prs.slides._sldIdLst  # 获取幻灯片ID列表
    slides = list(xml_slides)

    # 添加标题slide
    new_title_slide = copy_slide(prs, 0)
    for shape in new_title_slide.shapes:
        if not shape.has_text_frame:
            continue
        paragraphs = shape.text_frame.paragraphs
        # for i in range(1, len(paragraphs)):
        #     shape.text_frame.remove_paragraph(i)
        for run in paragraphs[0].runs:
            if "title" in run.text.lower():
                run.text = "大模型的前世今生"
            elif "speaker" in run.text.lower():
                run.text = "JOJO"

    new_index_slide = copy_slide(prs, 1)
    for shape in new_index_slide.shapes:
        if not shape.has_text_frame:
            continue
        paragraphs = shape.text_frame.paragraphs
        # for i in range(1, len(paragraphs)):
        #     shape.text_frame.remove_paragraph(i)
        for run in paragraphs[0].runs:
            if "title" in run.text.lower():
                run.text = "大模型的前世今生"
            elif "speaker" in run.text.lower():
                run.text = "JOJO"

    new_subtitle_slide = copy_slide(prs, 2)
    for shape in new_subtitle_slide.shapes:
        if not shape.has_text_frame:
            continue
        paragraphs = shape.text_frame.paragraphs
        # for i in range(1, len(paragraphs)):
        #     shape.text_frame.remove_paragraph(i)
        for run in paragraphs[0].runs:
            if "title" in run.text.lower():
                run.text = "Transformer介绍"
            elif "order" in run.text.lower():
                run.text = "01"

    new_content_slide = copy_slide(prs, 3)
    content_list = [
        ("优化沟通流程", "项目旨在改善公司内部沟通，提高工作效率。"),
        ("明确项目目标", "通过PPT汇报确保团队成员对目标与成果有共同理解。"),
        ("完成需求分析", "已进行需求分析并完成初步设计，确定项目实施路径。"),
        ("搭建协作平台", "团队协作平台已搭建完成，初步测试反馈良好。"),
        ("采用AI技术", "利用AI辅助设计提高PPT制作效率与质量，集成数据分析工具。")
    ]
    for shape in new_content_slide.shapes:
        if not shape.has_text_frame:
            continue
        paragraphs = shape.text_frame.paragraphs
        # for i in range(1, len(paragraphs)):
        #     shape.text_frame.remove_paragraph(i)
        for run in paragraphs[0].runs:
            subtitle_res = match_placeholder("subtitle", run.text)
            content_res = match_placeholder("content", run.text)
            if subtitle_res:
                run.text = content_list[content_res - 1][0]
            elif "title" in run.text.lower():
                run.text = "1.1 Transformer发展史"
            elif content_res:
                run.text = content_list[content_res - 1][1]

    for slide_id in xml_slides:
        prs.slides._sldIdLst.remove(slide_id)

    # 保存修改后的PPT文件
    prs.save(save_route)



    # 删除原本的slide
    for slide in slides:
        xml_slides.remove(slide)



if __name__ == '__main__':
    # 获取当前文件所在的路径
    dir_path = os.path.dirname(os.path.realpath(__file__))
    template_route = os.path.join(dir_path, "Template", "Design-qwen.pptx")
    save_route = os.path.join(dir_path, "GeneratedPresentations", "test.pptx")
    fill_ppt_content(template_route, save_route)
