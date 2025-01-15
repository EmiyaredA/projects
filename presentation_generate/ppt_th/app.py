import os
import random

from pptx import Presentation
import openai
from dataclasses import dataclass, field
from prompt import generate_ppt_page_prompt, generate_outline_system, generate_outline_user
import re
from collections import defaultdict
from bs4 import BeautifulSoup
from pprint import pprint

def __0_工具函数__():
    pass

def __0_1_模型对话__():
    pass

client = openai.OpenAI(
    base_url='https://platform.llmprovider.ai/v1',
    api_key='sk-Y7WuSEuYIchtepzPxcvIHxby7jkRHMxhDfWwSu4PE68eD70eC2C249Ae82221969B26a47Bb'
)
def chat_model(system_prompt, user_prompt):
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content":  user_prompt}
        ],
        temperature=0.5,
    )

    return response.choices[0].message.content

def __0_2_模型输出后处理__():
    pass

def parse_mermaid_from_md(md_text):
    try: 
        # 匹配Markdown代码块中的Mermaid图表
        mermaid_block_pattern = re.compile('```mermaid\n(.*?)\n```', re.DOTALL)

        # 提取Mermaid代码块
        mermaid_code_match = mermaid_block_pattern.search(md_text)
        mermaid_code = mermaid_code_match.group(1)
        # 解析Mermaid图表
        nodes, graph_structure = parse_mermaid_graph(mermaid_code)
        return nodes, graph_structure
    except Exception as e:
        raise f"解析mermaid失败: {e}"

# 解析Mermaid图表函数（之前定义的）
def parse_mermaid_graph(mermaid_code):
    node_pattern = re.compile(r'([A-Z0-9]+)\[([^\]]+)\]')
    connection_pattern = re.compile(r'([A-Z0-9]+) --> ([A-Z0-9]+)')

    nodes = {}
    connections = []

    for match in node_pattern.finditer(mermaid_code):
        node_id, node_label = match.groups()
        nodes[node_id] = node_label

    for match in connection_pattern.finditer(mermaid_code):
        from_node, to_node = match.groups()
        connections.append((from_node, to_node))

    graph_structure = {}
    for from_node, to_node in connections:
        if from_node not in graph_structure:
            graph_structure[from_node] = []
        graph_structure[from_node].append(to_node)

    return nodes, graph_structure

def parsed_html_content(html_content):
    # Parse HTML using BeautifulSoup
    soup = BeautifulSoup(html_content, 'html.parser')

    # Find all the <h1> and <p> elements
    titles = soup.find_all('h1')
    contents = soup.find_all('p')

    # Create a list of dictionaries containing title and content
    result = []
    for title, content in zip(titles, contents):
        result.append({
            'title': title.get_text(),
            'content': content.get_text()
        })

    # Output the result
    return result

def find_root(nodes):
    """
    自动确定根节点
    :param data: 原始数据字典
    :return: 根节点
    """
    # 获取所有节点
    all_nodes = set(nodes.keys())
    # 获取所有子节点
    child_nodes = set(child for children in nodes.values() for child in children)
    # 根节点是出现在数据中但不在任何子节点列表中的节点
    root_nodes = all_nodes - child_nodes
    return root_nodes.pop() if root_nodes else None

def build_tree(data, root):
    """
    根据原始数据字典构建树形结构字典
    :param data: 原始数据字典
    :param root: 当前要处理的根节点
    :return: 一个树形结构字典
    """
    tree = {}
    # 获取当前节点的子节点列表
    if root in data:
        for child in data[root]:
            # 递归构建子节点的树结构
            tree[child] = build_tree(data, child)
    return tree

def __1_业务函数__():
    pass

@dataclass
class OutlineConfig:
    topic: str = ""
    speaker: str = ""
    content: str = ""
    language: str = "中文"
    slide_count: int = 10

pair_count_list = []
def get_random_pair_count():
    if not pair_count_list:
        # 如果列表为空，重新填充列表
        pair_count_list.extend([1, 1, 2, 2, 3, 3, 4, 4, 5, 6, 8])  # 根据实际需求填充
    item = random.choice(pair_count_list)  # 随机选择一个元素
    pair_count_list.remove(item)  # 从列表中移除该元素
    return item

def set_content(slide, text, type):
    match type:
        case "title":
            title = slide.shapes.title
            title.text = text
        case "content":
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = text
        case _:
            raise ValueError("设置ppt内容阶段，slide中元素的类型传参错误")

def generate_sub_ppt_pages(prs, nodes, sub_tree, curr_node, path):

    if not path:
        raise "生成ppt子页的路径不能为空"
        # 将当前节点加入路径
    path.append(nodes[curr_node])

    if not sub_tree:
        # 如果当前节点没有子节点，说明是叶子节点，开始生成对应的ppt子页
        pair_count = get_random_pair_count()
        generate_ppt_page_prompt_str = generate_ppt_page_prompt(pair_count)
        answer = chat_model(
            system_prompt=generate_ppt_page_prompt_str,
            user_prompt=str(path)
        )
        parsed_answer_list = parsed_html_content(answer)
        slide = prs.slides.add_slide(prs.slide_layouts[7])
        content = '\n\n'.join([
            f"{parsed_answer_item['title']}：\n{parsed_answer_item['content']}"
            for parsed_answer_item in parsed_answer_list
        ])
        print(f"路径：{str(path)}\n内容：{content}\n\n")
        set_content(slide, content, "content")
        return

    # 遍历子节点
    for child_key in sub_tree:
        generate_sub_ppt_pages(prs, nodes, sub_tree[child_key], child_key, path)
    return

def generate_ppt(nodes, graph_structure, template_path, config):

    # 自动确定根节点
    root_node = find_root(graph_structure)
    if root_node:
        tree_structure = {root_node: build_tree(graph_structure, root_node)}
        pprint(tree_structure)
        print()
    else:
        raise "生成的大纲树没有根节点"

    title = nodes[root_node]
    # 加载PPT模板
    prs = Presentation(template_path)
    # 创建标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_content(slide, title, "title")
    set_content(slide, f"演讲人：{config.speaker}", "content")
    # 创建目录页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_content(slide, "目录", "title")
    dir_list = [nodes[child] for child in tree_structure[root_node]]
    set_content(slide, "\n".join(dir_list), "content")
    # 遍历树结构，为每个子节点创建幻灯片
    for child_node in tree_structure[root_node]:
        # 创建子标题页
        slide = prs.slides.add_slide(prs.slide_layouts[2])
        sub_title = nodes[child_node]
        set_content(slide, sub_title, "title")
        # 创建子内容页
        generate_sub_ppt_pages(prs, nodes, tree_structure[root_node][child_node], child_node, [title, sub_title])

    prs.save(f'GeneratedPresentations/{title}.pptx')
    # file_path = f"GeneratedPresentations/{title}.pptx"

def generate_outline(config: OutlineConfig = OutlineConfig()):
    generate_outline_system_str = generate_outline_system.render()
    generate_outline_user_str = generate_outline_user.render(
        topic = config.topic,
        # speaker = config.speaker,
        content = config.content,
        language = config.language,
        slide_count = config.slide_count
    )
    answer = chat_model(system_prompt=generate_outline_system_str, user_prompt=generate_outline_user_str)
    print(answer)
    print()
    nodes, graph_structure = parse_mermaid_from_md(answer)
    if nodes is None or graph_structure is None:
        raise ValueError("大纲生成失败")
    return nodes, graph_structure

if __name__ == '__main__':
    current_dir = os.path.dirname(__file__)  # 获取当前脚本所在的目录
    ppt_template_path = os.path.join(current_dir, 'Template', 'Design-1.pptx')
    # ppt_template_path = os.path.join(current_dir, 'Template', 'Design-qwen.pptx')

    outlineConfig = OutlineConfig(
        topic="python教程",
        speaker="jojo",
        content="详细地介绍以下python的使用",
        language="中文",
        slide_count=30
    )
    nodes, graph_structure = generate_outline(outlineConfig)
    generate_ppt(nodes, graph_structure, ppt_template_path, outlineConfig)